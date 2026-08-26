#!/usr/bin/env python3
"""Build the integrated Cluster2 steal-buffer polarity AER presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation" / "cluster2_digital_first_round_20260828.pptx"

BG = RGBColor(7, 22, 39)
CARD = RGBColor(18, 48, 76)
CARD_2 = RGBColor(25, 62, 92)
WHITE = RGBColor(246, 249, 252)
MUTED = RGBColor(172, 194, 214)
CYAN = RGBColor(42, 211, 205)
GREEN = RGBColor(55, 205, 133)
BLUE = RGBColor(71, 147, 255)
ORANGE = RGBColor(255, 183, 70)
RED = RGBColor(255, 99, 118)
PURPLE = RGBColor(178, 126, 255)
FONT = "Malgun Gothic"


def shape(slide, x, y, w, h, fill, rounded=True, line=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    obj = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    obj.fill.solid()
    obj.fill.fore_color.rgb = fill
    if line is None:
        obj.line.fill.background()
    else:
        obj.line.color.rgb = line
        obj.line.width = Pt(1.5)
    return obj


def text(slide, value, x, y, w, h, size=24, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = valign
    frame.margin_left = frame.margin_right = Inches(0.02)
    frame.margin_top = frame.margin_bottom = Inches(0.01)
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def header(prs, section, title, number, source):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    shape(slide, 0, 0, 0.14, 7.5, CYAN, False)
    shape(slide, 0.55, 0.26, 1.40, 0.34, CARD_2)
    text(slide, section, 0.60, 0.30, 1.30, 0.22, 11, CYAN, True, PP_ALIGN.CENTER)
    text(slide, title, 2.18, 0.18, 10.45, 0.58, 27, WHITE, True)
    shape(slide, 0.55, 0.88, 12.18, 0.025, CYAN, False)
    text(slide, source, 0.60, 7.10, 11.30, 0.20, 8.5, MUTED)
    text(slide, f"{number:02d}", 12.02, 7.08, 0.60, 0.20, 10, CYAN, True, PP_ALIGN.RIGHT)
    return slide


def arrow(slide, x, y, color=CYAN, size=28, w=0.55):
    text(slide, "→", x, y, w, 0.42, size, color, True, PP_ALIGN.CENTER)


def takeaway(slide, value, color=GREEN):
    shape(slide, 0.78, 6.20, 11.78, 0.62, CARD)
    text(slide, value, 1.02, 6.32, 11.30, 0.34, 17, color, True, PP_ALIGN.CENTER)


def metric(slide, x, y, w, label, value, accent, note=""):
    shape(slide, x, y, w, 1.15, CARD, True, accent)
    text(slide, label, x + 0.18, y + 0.12, w - 0.36, 0.24, 11, MUTED, True)
    text(slide, value, x + 0.18, y + 0.40, w - 0.36, 0.42, 24, accent, True)
    if note:
        text(slide, note, x + 0.18, y + 0.86, w - 0.36, 0.19, 10, MUTED)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. The complete project in one sentence.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    shape(slide, 0, 0, 0.18, 7.5, CYAN, False)
    text(slide, "병목을 줄이고 의미를 보존하는 AER", 0.84, 0.68, 11.8, 0.76, 38, WHITE, True)
    text(slide, "Cluster2 Steal-Buffer Polarity AER", 0.86, 1.55, 10.8, 0.38, 21, CYAN, True)
    stages = [
        (0.90, RED, "AER 병목", "직렬 출력\n재발생·불균형"),
        (4.43, CYAN, "제안 구조", "bitmap · FIFO\nsteal · polarity"),
        (8.37, GREEN, "검증 결과", "full50 loss 0.47%\nUZH 8,503 / 8,503"),
    ]
    for i, (x, accent, tag, body) in enumerate(stages):
        shape(slide, x, 2.45, 3.15, 2.05, CARD, True, accent)
        text(slide, tag, x + 0.25, 2.72, 2.65, 0.26, 13, accent, True, PP_ALIGN.CENTER)
        text(slide, body, x + 0.25, 3.18, 2.65, 0.82, 22, WHITE, True, PP_ALIGN.CENTER)
        if i < 2:
            arrow(slide, x + 3.12, 3.20, CYAN, 32, 0.40)
    shape(slide, 0.90, 5.18, 10.62, 0.52, CARD_2)
    text(slide, "문제 정의  →  구조 설계  →  정량 개선  →  RTL/PPA  →  CAV 확장", 1.12, 5.31, 10.18, 0.24, 16, MUTED, True, PP_ALIGN.CENTER)
    text(slide, "2026-08-28  |  Xcelium · Genus · Innovus", 0.86, 6.86, 11.0, 0.24, 10, MUTED)

    # 2. Translate all six literature problems into project requirements.
    slide = header(prs, "01 문제", "Ryu의 문제 흐름을 여섯 설계 요구사항으로 분해했다", 2,
                   "Source: team decomposition of Ryu et al., CVPRW 2019 · direct = scoped mitigation")
    rows = [
        ("① 주소 overhead", "bitmap 집약", "repeat-flag 후속", ORANGE, "부분"),
        ("② bandwidth", "2 row-bitmap lanes", "최대 8 events/cycle", GREEN, "직접"),
        ("③ arbitration latency", "병렬 retire + depth-2", "대기·재발생 완화", GREEN, "직접"),
        ("④ unfair arbitration", "rotating RR + steal", "fairness 수치는 미측정", PURPLE, "구조"),
        ("⑤ timestamp", "time sidecar only", "DUT payload 통합은 후속", RED, "HOLD"),
        ("⑥ motion artifact", "sensor/system 대응", "이번 RTL에서는 미검증", RED, "HOLD"),
    ]
    for i, (problem, action, result, accent, status) in enumerate(rows):
        x = 0.78 + (i % 2) * 6.02
        y = 1.22 + (i // 2) * 1.48
        shape(slide, x, y, 5.72, 1.16, CARD, True, accent)
        shape(slide, x + 0.20, y + 0.18, 0.75, 0.30, accent)
        text(slide, status, x + 0.25, y + 0.21, 0.65, 0.20, 10, BG, True, PP_ALIGN.CENTER)
        text(slide, problem, x + 1.10, y + 0.18, 2.15, 0.25, 13, accent, True)
        text(slide, action, x + 1.10, y + 0.55, 2.15, 0.28, 14, WHITE, True)
        text(slide, result, x + 3.18, y + 0.53, 2.22, 0.32, 12, MUTED, False, PP_ALIGN.RIGHT)
    takeaway(slide, "②③은 직접 완화 · ①④는 부분/구조 대응 · ⑤⑥은 이번 RTL의 HOLD 범위다", ORANGE)

    # 3. Measured baseline mechanism.
    slide = header(prs, "01 문제", "scalar 출력은 동시 입력을 source-local loss로 바꾼다", 3,
                   "Source: recovered scalar Fovea Xcelium full50 diagnostics")
    text(slide, "동시 입력", 0.85, 1.24, 2.10, 0.25, 13, BLUE, True, PP_ALIGN.CENTER)
    for r in range(3):
        for c in range(4):
            shape(slide, 0.94 + c * 0.46, 1.68 + r * 0.46, 0.33, 0.33, CARD_2, True, BLUE)
    arrow(slide, 3.02, 2.17, ORANGE, 34, 0.70)
    shape(slide, 3.92, 1.40, 3.24, 2.10, CARD)
    text(slide, "공유 arbitration", 4.20, 1.75, 2.68, 0.32, 18, WHITE, True, PP_ALIGN.CENTER)
    text(slide, "여러 event가 하나의\n출력 순서를 기다림", 4.34, 2.30, 2.42, 0.70, 15, MUTED, False, PP_ALIGN.CENTER)
    arrow(slide, 7.35, 2.17, RED, 34, 0.70)
    shape(slide, 8.24, 1.40, 4.05, 2.10, CARD, True, RED)
    text(slide, "1 EVENT / CYCLE", 8.56, 1.82, 3.42, 0.40, 23, RED, True, PP_ALIGN.CENTER)
    text(slide, "대기 누적  →  local full  →  overrun", 8.48, 2.57, 3.56, 0.34, 14, WHITE, True, PP_ALIGN.CENTER)
    metric(slide, 0.92, 4.12, 3.55, "ACCEPTED", "78,229 / 106,416", ORANGE, "73.51%")
    metric(slide, 4.90, 4.12, 3.05, "OVERRUN", "28,187", RED, "full50 baseline")
    metric(slide, 8.38, 4.12, 3.90, "THROUGHPUT", "0.673901 EPC", RED, "fixed-window")
    takeaway(slide, "따라서 출력 폭, 재발생 흡수, traffic 균형을 동시에 풀어야 한다", ORANGE)

    # 4. The complete proposed architecture.
    slide = header(prs, "02 구조", "네 메커니즘이 하나의 retire 경로에서 함께 동작한다", 4,
                   "Source: final polarity RTL SHA-256 20d601a9…")
    blocks = [
        (0.62, 1.55, 2.05, 2.50, BLUE, "16 SOURCES", "event + polarity"),
        (3.10, 1.55, 2.40, 2.50, ORANGE, "16 × DEPTH-2", "pending count\n+ 2 polarity slots"),
        (5.95, 1.15, 2.55, 1.55, CYAN, "CENTER LANE", "row bitmap"),
        (5.95, 3.05, 2.55, 1.55, PURPLE, "PERIPH. LANE", "row bitmap"),
        (9.00, 1.55, 3.15, 2.50, GREEN, "RETIRE", "row + col_mask\n+ pol_mask"),
    ]
    for x, y, w, h, accent, tag, body in blocks:
        shape(slide, x, y, w, h, CARD, True, accent)
        text(slide, tag, x + 0.18, y + 0.28, w - 0.36, 0.26, 12, accent, True, PP_ALIGN.CENTER)
        text(slide, body, x + 0.20, y + 0.88, w - 0.40, 0.85, 18, WHITE, True, PP_ALIGN.CENTER)
    arrow(slide, 2.69, 2.52, CYAN, 28, 0.38)
    arrow(slide, 5.52, 2.52, CYAN, 28, 0.38)
    arrow(slide, 8.54, 2.52, GREEN, 28, 0.38)
    shape(slide, 5.75, 4.88, 2.95, 0.46, CARD_2, True, PURPLE)
    text(slide, "traffic 쏠림 시 lane steal", 5.92, 4.98, 2.62, 0.24, 12, PURPLE, True, PP_ALIGN.CENTER)
    shape(slide, 9.34, 4.88, 2.45, 0.46, GREEN)
    text(slide, "최대 8 events/cycle", 9.48, 4.98, 2.18, 0.24, 12, BG, True, PP_ALIGN.CENTER)
    takeaway(slide, "Cluster2의 병렬 row 전송 위에 depth-2·steal·polarity lockstep을 통합했다", GREEN)

    # 5. Bitmap and depth-2 solve different causes.
    slide = header(prs, "02 구조", "bitmap은 출력 폭, depth-2는 source 재발생을 푼다", 5,
                   "Source: final RTL contract · full50 bottleneck diagnosis")
    shape(slide, 0.78, 1.25, 5.72, 4.55, CARD, True, CYAN)
    text(slide, "IDEA 1  ROW BITMAP", 1.08, 1.56, 5.12, 0.28, 14, CYAN, True, PP_ALIGN.CENTER)
    for c in range(4):
        shape(slide, 1.25 + c * 0.80, 2.22, 0.52, 0.52, CARD_2, True, CYAN)
    text(slide, "같은 row의 4 columns", 1.12, 2.96, 4.95, 0.25, 14, MUTED, True, PP_ALIGN.CENTER)
    arrow(slide, 2.92, 3.34, CYAN, 30, 0.55)
    shape(slide, 1.48, 3.92, 4.30, 0.88, CARD_2, True, CYAN)
    text(slide, "col_mask = 4'b1111", 1.78, 4.16, 3.70, 0.34, 20, WHITE, True, PP_ALIGN.CENTER)
    text(slide, "1개 주소가 최대 4 events 표현", 1.30, 5.16, 4.70, 0.24, 13, CYAN, True, PP_ALIGN.CENTER)
    shape(slide, 6.82, 1.25, 5.72, 4.55, CARD, True, ORANGE)
    text(slide, "IDEA 2  DEPTH-2 / SOURCE", 7.12, 1.56, 5.12, 0.28, 14, ORANGE, True, PP_ALIGN.CENTER)
    text(slide, "t", 7.28, 2.10, 0.36, 0.22, 11, MUTED, True)
    for i, label in enumerate(["event A", "event B", "grant A"]):
        x = 7.75 + i * 1.40
        shape(slide, x, 2.04, 1.05, 0.58, CARD_2, True, ORANGE if i < 2 else GREEN)
        text(slide, label, x + 0.08, 2.19, 0.89, 0.23, 11, WHITE, True, PP_ALIGN.CENTER)
    shape(slide, 8.08, 3.15, 3.28, 1.15, CARD_2, True, ORANGE)
    text(slide, "slot 0: A  |  slot 1: B", 8.30, 3.52, 2.84, 0.30, 17, WHITE, True, PP_ALIGN.CENTER)
    text(slide, "두 번째 event를 overrun 대신 저장", 7.22, 4.75, 4.92, 0.28, 14, ORANGE, True, PP_ALIGN.CENTER)
    takeaway(slide, "병렬 표현과 local buffering을 분리했기 때문에 두 병목을 동시에 완화한다", GREEN)

    # 6. Steal and polarity complete the design.
    slide = header(prs, "02 구조", "lane steal은 불균형을, lockstep은 event 의미를 지킨다", 6,
                   "Source: final polarity RTL · native observational contract")
    shape(slide, 0.78, 1.25, 5.72, 4.55, CARD, True, PURPLE)
    text(slide, "IDEA 3  CONDITIONAL STEAL", 1.08, 1.56, 5.12, 0.28, 14, PURPLE, True, PP_ALIGN.CENTER)
    shape(slide, 1.30, 2.18, 1.82, 1.50, CARD_2, True, CYAN)
    text(slide, "CENTER", 1.55, 2.44, 1.32, 0.22, 12, CYAN, True, PP_ALIGN.CENTER)
    text(slide, "busy", 1.55, 2.93, 1.32, 0.32, 19, WHITE, True, PP_ALIGN.CENTER)
    shape(slide, 4.15, 2.18, 1.82, 1.50, CARD_2, True, MUTED)
    text(slide, "PERIPH.", 4.40, 2.44, 1.32, 0.22, 12, MUTED, True, PP_ALIGN.CENTER)
    text(slide, "idle", 4.40, 2.93, 1.32, 0.32, 19, MUTED, True, PP_ALIGN.CENTER)
    text(slide, "← steal\ncapacity", 3.02, 2.66, 1.16, 0.50, 12, PURPLE, True, PP_ALIGN.CENTER)
    text(slide, "idle class의 처리력을 재사용\n(global fairness 수치는 별도)", 1.20, 4.43, 4.90, 0.62, 13, PURPLE, True, PP_ALIGN.CENTER)
    shape(slide, 6.82, 1.25, 5.72, 4.55, CARD, True, BLUE)
    text(slide, "IDEA 4  POLARITY LOCKSTEP", 7.12, 1.56, 5.12, 0.28, 14, BLUE, True, PP_ALIGN.CENTER)
    bundles = [("A", "row2 col1", "pol 0"), ("B", "row2 col2", "pol 1")]
    for i, (name, addr, pol) in enumerate(bundles):
        y = 2.16 + i * 1.15
        shape(slide, 7.45, y, 4.46, 0.82, CARD_2, True, BLUE)
        text(slide, name, 7.68, y + 0.22, 0.40, 0.25, 14, BLUE, True)
        text(slide, addr, 8.20, y + 0.22, 1.72, 0.25, 13, WHITE, True)
        text(slide, pol, 10.24, y + 0.22, 1.35, 0.25, 13, BLUE, True, PP_ALIGN.RIGHT)
    text(slide, "주소·polarity를 같은 slot에서 push/pop", 7.18, 4.70, 4.98, 0.28, 14, BLUE, True, PP_ALIGN.CENTER)
    takeaway(slide, "steal_buf는 처리력 낭비와 재발생 손실을 잡고, polarity FIFO는 의미를 보존한다", GREEN)

    # 7. Address-overhead extensions are presented as a coherent next layer.
    slide = header(prs, "02 구조", "주소 부호화: repeat-flag 통합은 아직 HOLD", 7,
                   "Source: row-trim/repeat-flag studies · address-only integration boundary")
    codec_cards = [
        (0.78, 3.60, CYAN, "FINAL POLARITY RTL", "ROW + MASK", "row · col_mask · pol_mask", "제출된 polarity RTL의 현재 encoding", MUTED),
        (4.78, 3.48, MUTED, "ADDRESS-ONLY STUDY", "ROW-TRIM", "−14.29%", "기본 Cluster2 전용\nsteal_buf row에는 적용 불가", RED),
        (8.66, 3.88, ORANGE, "ADDRESS-ONLY STUDY", "REPEAT-FLAG", "−15.61%", "steal_buf address 실험\npolarity grammar 재검증", ORANGE),
    ]
    for x, w, accent, badge, name, value, note, note_color in codec_cards:
        shape(slide, x, 1.35, w, 4.42, CARD, True, accent)
        text(slide, badge, x + 0.22, 1.66, w - 0.44, 0.30, 16, accent, True, PP_ALIGN.CENTER)
        text(slide, name, x + 0.22, 2.18, w - 0.44, 0.38, 22, WHITE, True, PP_ALIGN.CENTER)
        value_size = 34 if value.startswith("−") else 18
        text(slide, value, x + 0.22, 2.86, w - 0.44, 0.62, value_size, accent, True, PP_ALIGN.CENTER)
        text(slide, note, x + 0.25, 4.12, w - 0.50, 0.78, 16, note_color, True, PP_ALIGN.CENTER)
    takeaway(slide, "최종 제출: ROW + MASK  |  repeat-flag: final-polarity 재검증 전까지 2차 후보", ORANGE)

    # 8. One integrated quantitative evolution.
    slide = header(prs, "03 결과", "full50 source-overrun: 12,259 → 502 (−95.9%)", 8,
                   "Source: recovered scalar/base + upstream steal_buf full50 · same 106,416-event denominator")
    shape(slide, 0.78, 1.04, 11.78, 0.48, CARD_2)
    text(slide, "동일 full50 50 traces · 106,416 generated events / design", 1.02, 1.16, 11.30, 0.24, 15, CYAN, True, PP_ALIGN.CENTER)
    data = [
        ("SCALAR FOVEA", 26.49, "28,187", RED),
        ("BASE CLUSTER2", 11.52, "12,259", ORANGE),
        ("STEAL_BUF", 0.47, "502", GREEN),
    ]
    max_h = 2.75
    for i, (label, pct, count, accent) in enumerate(data):
        x = 0.98 + i * 2.98
        h = max(0.10, max_h * pct / 26.49)
        shape(slide, x + 0.52, 4.78 - h, 1.48, h, accent, False)
        text(slide, f"{pct:.2f}%", x + 0.12, 1.72 if i == 2 else 4.38 - h, 2.28, 0.42, 24, accent, True, PP_ALIGN.CENTER)
        text(slide, label, x, 5.00, 2.52, 0.30, 13, WHITE, True, PP_ALIGN.CENTER)
        text(slide, f"loss {count}", x, 5.40, 2.52, 0.24, 12, MUTED, False, PP_ALIGN.CENTER)
    shape(slide, 9.78, 1.55, 2.55, 1.50, GREEN)
    text(slide, "−95.9%", 10.00, 1.82, 2.11, 0.46, 27, BG, True, PP_ALIGN.CENTER)
    text(slide, "남은 loss = 1/24.42", 9.92, 2.40, 2.28, 0.26, 12, BG, True, PP_ALIGN.CENTER)
    shape(slide, 9.58, 3.40, 2.95, 1.18, CARD, True, ORANGE)
    text(slide, "48 / 50 traces", 9.80, 3.63, 2.50, 0.28, 16, GREEN, True, PP_ALIGN.CENTER)
    text(slide, "overrun 0 · worst 5.39%", 9.74, 4.04, 2.62, 0.26, 11, MUTED, True, PP_ALIGN.CENTER)
    takeaway(slide, "architecture-family full50 결과이며, 다음 UZH polarity-v1 검증과는 별도 population이다", ORANGE)

    # 9. Actual simulator and cycle evidence.
    slide = header(prs, "03 결과", "실제 TB: 입력부터 cycle ledger까지 독립 검증", 9,
                   "Source: Xcelium 23.09-s013 · polarity-v1 native observational evidence")
    stages = [
        (0.72, BLUE, "UZH TRACE", "8,503 events"),
        (3.48, CYAN, "RTL + TB", "Xcelium\n23.09-s013"),
        (6.58, PURPLE, "CYCLE LEDGER", "independent check"),
        (9.72, GREEN, "RESULT", "8,503 / 8,503"),
    ]
    for i, (x, accent, label, value) in enumerate(stages):
        shape(slide, x, 1.34, 2.42, 1.25, CARD, True, accent)
        text(slide, label, x + 0.18, 1.55, 2.06, 0.25, 13, accent, True, PP_ALIGN.CENTER)
        text(slide, value, x + 0.14, 1.87, 2.14, 0.52, 16, WHITE, True, PP_ALIGN.CENTER)
        if i < 3:
            arrow(slide, x + 2.48, 1.75, CYAN, 24, 0.36)
    text(slide, "TB · redred_cluster2_polarity_v1_native_observational_tb.sv",
         1.20, 2.67, 10.90, 0.24, 10, MUTED, True, PP_ALIGN.CENTER)
    shape(slide, 0.90, 3.02, 11.55, 1.34, CARD, True, CYAN)
    text(slide, "REAL RETIRE · CYCLE 4162", 1.18, 3.27, 2.70, 0.25, 12, CYAN, True)
    text(slide, "lane0  row2  col=0x6  pol=0x0", 4.05, 3.23, 3.42, 0.30, 16, WHITE, True, PP_ALIGN.CENTER)
    text(slide, "lane1  row0  col=0x1  pol=0x1", 7.83, 3.23, 3.42, 0.30, 16, WHITE, True, PP_ALIGN.CENTER)
    text(slide, "col_mask popcount 2 + 1  →  두 lane에서 3 events 동시 retire", 3.02, 3.82, 7.30, 0.25, 13, GREEN, True, PP_ALIGN.CENTER)
    checks = [("overrun", "0"), ("phantom", "0"), ("duplicate", "0"), ("pol mismatch", "0")]
    for i, (label, value) in enumerate(checks):
        x = 1.08 + i * 2.84
        shape(slide, x, 4.77, 2.46, 0.82, CARD_2, True, GREEN)
        text(slide, label, x + 0.12, 4.94, 1.56, 0.24, 12, MUTED, True)
        text(slide, value, x + 1.65, 4.91, 0.55, 0.30, 19, GREEN, True, PP_ALIGN.RIGHT)
    takeaway(slide, "UZH shapes_rotation 4×4 patch · 1 ms bins · selected-event polarity sequence · drain empty", GREEN)

    # 10. Physical results.
    slide = header(prs, "04 구현", "최종 polarity-v1 RTL: 3.5 ns post-route setup·hold PASS", 10,
                   "Source: Genus/Innovus · GPDK045 slow 0.9 V/125 °C · discrete per-target sweep")
    shape(slide, 0.78, 1.20, 4.10, 1.38, CARD, True, GREEN)
    text(slide, "POST-ROUTE TIMING TARGET", 0.98, 1.36, 2.70, 0.22, 10, MUTED, True)
    shape(slide, 3.78, 1.33, 0.78, 0.30, GREEN)
    text(slide, "PASS", 3.84, 1.39, 0.66, 0.18, 10, BG, True, PP_ALIGN.CENTER)
    text(slide, "3.5 ns · 285.714 MHz", 0.98, 1.68, 3.68, 0.42, 25, GREEN, True)
    text(slide, "setup WNS +0.454 ns · hold WNS +0.167 ns", 0.98, 2.20, 3.68, 0.20, 10, MUTED)
    shape(slide, 5.18, 1.20, 3.38, 1.38, CARD, True, CYAN)
    text(slide, "AREA · RAW REPORT UNIT", 5.38, 1.36, 2.98, 0.22, 10, MUTED, True)
    text(slide, "SYNTHESIS", 5.38, 1.72, 1.08, 0.22, 10, MUTED, True)
    text(slide, "1156.644", 6.48, 1.67, 1.68, 0.28, 19, CYAN, True, PP_ALIGN.RIGHT)
    shape(slide, 5.38, 2.02, 2.78, 0.01, CARD_2, False)
    text(slide, "P&R", 5.38, 2.14, 1.08, 0.22, 10, MUTED, True)
    text(slide, "1254.114", 6.48, 2.09, 1.68, 0.28, 19, WHITE, True, PP_ALIGN.RIGHT)
    shape(slide, 8.86, 1.20, 3.46, 1.38, CARD, True, ORANGE)
    text(slide, "POST-ROUTE POWER · VECTORLESS", 9.06, 1.36, 3.06, 0.22, 10, MUTED, True)
    text(slide, "0.107389 mW", 9.06, 1.68, 3.06, 0.42, 24, ORANGE, True)
    text(slide, "default activity 0.2 · not workload power", 9.06, 2.20, 3.06, 0.20, 10, MUTED)
    text(slide, "DISCRETE POST-ROUTE SETUP WNS SWEEP", 1.02, 2.78, 5.00, 0.22, 10, MUTED, True)
    shape(slide, 1.66, 3.28, 5.96, 0.08, GREEN, False)
    shape(slide, 7.62, 3.28, 2.98, 0.08, CARD_2, False)
    points = [
        (1.66, "222.222", "+1.349", GREEN),
        (4.64, "250.000", "+0.849", GREEN),
        (7.62, "285.714", "+0.454", GREEN),
        (10.60, "333.333", "−0.004", RED),
    ]
    for center, freq, slack, accent in points:
        shape(slide, center - 0.23, 3.05, 0.46, 0.46, accent)
        text(slide, f"{freq} MHz", center - 0.71, 3.72, 1.42, 0.26, 13, accent, True, PP_ALIGN.CENTER)
        text(slide, f"WNS {slack} ns", center - 0.77, 4.10, 1.54, 0.22, 10, MUTED, False, PP_ALIGN.CENTER)
    shape(slide, 0.98, 4.62, 7.16, 0.78, CARD, True, CYAN)
    text(slide, "PER-TARGET FLOW", 1.20, 4.75, 2.10, 0.20, 10, MUTED, True)
    text(slide, "Genus: generic → map → opt  |  Innovus: place → CTS → route", 1.20, 5.02, 6.72, 0.26, 13, WHITE, True, PP_ALIGN.CENTER)
    shape(slide, 8.48, 4.62, 3.62, 0.78, CARD, True, GREEN)
    text(slide, "INTERNAL TOOL CHECKS", 8.68, 4.75, 3.22, 0.20, 10, MUTED, True)
    text(slide, "DRC 0 · antenna 0", 8.68, 5.02, 3.22, 0.26, 15, GREEN, True, PP_ALIGN.CENTER)
    shape(slide, 0.98, 5.52, 11.12, 0.48, CARD, True, ORANGE)
    text(slide, "LIMITS", 1.16, 5.65, 0.74, 0.20, 10, ORANGE, True)
    text(slide, "single slow view · Non-OCV · SI off · no SPEF/RCDB · vectorless power · internal checks only",
         1.98, 5.64, 9.90, 0.22, 11, WHITE, True)
    takeaway(slide, "GO: 3.5 ns reported setup·hold PASS  |  HOLD: exact Fmax·activity power·signoff", ORANGE)

    # 11. CAV extension with proven and next clearly separated.
    slide = header(prs, "05 확장", "별도 software path에서 CAV 좌표계 전수 분기를 검증했다", 11,
                   "Source: sealed legacy address-only official UZH→CAV result")
    shape(slide, 0.78, 1.14, 11.78, 0.52, CARD_2)
    text(slide, "SEPARATE SOFTWARE TRACK · legacy address-only", 1.02, 1.27, 11.30, 0.24, 14, PURPLE, True, PP_ALIGN.CENTER)
    stages = [
        (0.82, BLUE, "EVENTS + POSE", "8,503 INPUT"),
        (3.56, GREEN, "TB-SIDE JOIN", "8,503 / 8,503"),
        (6.52, PURPLE, "TIME SEMANTICS", "OCCURRENCE TIME"),
    ]
    for i, (x, accent, label, value) in enumerate(stages):
        shape(slide, x, 1.94, 2.34, 1.58, CARD, True, accent)
        text(slide, label, x + 0.17, 2.18, 2.00, 0.27, 14, accent, True, PP_ALIGN.CENTER)
        text(slide, value, x + 0.12, 2.66, 2.10, 0.52, 20, WHITE, True, PP_ALIGN.CENTER)
        if i < 2:
            arrow(slide, x + 2.40, 2.48, CYAN, 24, 0.40)
    shape(slide, 8.88, 2.72, 0.48, 0.04, CYAN, False)
    shape(slide, 9.32, 2.24, 0.04, 1.02, CYAN, False)
    shape(slide, 9.32, 2.24, 0.34, 0.04, CYAN, False)
    shape(slide, 9.32, 3.22, 0.34, 0.04, ORANGE, False)
    shape(slide, 9.66, 1.76, 2.55, 1.12, CARD, True, CYAN)
    text(slide, "WORLD", 9.88, 1.99, 2.12, 0.24, 14, CYAN, True, PP_ALIGN.CENTER)
    text(slide, "8,420", 9.88, 2.36, 2.12, 0.34, 24, WHITE, True, PP_ALIGN.CENTER)
    shape(slide, 9.66, 3.02, 2.55, 1.12, CARD, True, ORANGE)
    text(slide, "SENSOR_FIXED", 9.84, 3.25, 2.20, 0.24, 14, ORANGE, True, PP_ALIGN.CENTER)
    text(slide, "83", 9.88, 3.62, 2.12, 0.34, 24, WHITE, True, PP_ALIGN.CENTER)
    shape(slide, 0.98, 4.52, 5.60, 0.82, CARD)
    text(slide, "WORLD → 821 cells · ID는 RTL payload 아님", 1.20, 4.77, 5.16, 0.32, 16, PURPLE, True, PP_ALIGN.CENTER)
    shape(slide, 6.74, 4.52, 5.46, 0.82, CARD, True, ORANGE)
    text(slide, "NEXT · polarity replay → CAV RTL/PPA", 6.98, 4.77, 4.98, 0.32, 16, ORANGE, True, PP_ALIGN.CENTER)
    takeaway(slide, "확장성은 software 기능 경로로 입증했고, polarity→CAV RTL/PPA는 2차 과제다", ORANGE)

    # 12. Close with integrated conclusions and roadmap.
    slide = header(prs, "SUMMARY", "결론: RTL/PPA와 CAV software 경로를 각각 검증했다", 12,
                   "Source bundle: receipts · raw reports · SHA256SUMS · PROVENANCE.json")
    claims = [
        (0.78, CYAN, "DESIGN", "bitmap · FIFO", "steal · polarity"),
        (4.72, GREEN, "DIRECTIONAL EFFECT", "11.52% → 0.47%", "95.9% lower"),
        (8.66, BLUE, "PROVED FINAL RTL", "8,503 / 8,503", "285.714 MHz\ntested PASS"),
    ]
    for x, accent, tag, line1, line2 in claims:
        shape(slide, x, 1.38, 3.52, 2.22, CARD, True, accent)
        text(slide, tag, x + 0.25, 1.64, 3.02, 0.28, 15, accent, True, PP_ALIGN.CENTER)
        text(slide, line1, x + 0.22, 2.12, 3.08, 0.40, 23, WHITE, True, PP_ALIGN.CENTER)
        text(slide, line2, x + 0.20, 2.70, 3.12, 0.56, 17, accent, True, PP_ALIGN.CENTER)
    shape(slide, 0.78, 4.05, 11.40, 1.45, CARD)
    tracks = [
        (0.98, CYAN, "LINK", "repeat-flag 통합\nlink-bit 재측정"),
        (4.55, PURPLE, "CAV", "polarity full replay\n→ CAV RTL/PPA"),
        (8.12, ORANGE, "PPA", "activity power\nexact Fmax 탐색"),
    ]
    for x, accent, tag, body in tracks:
        shape(slide, x, 4.25, 3.16, 1.02, CARD_2, True, accent)
        text(slide, tag, x + 0.20, 4.43, 0.72, 0.24, 14, accent, True)
        text(slide, body, x + 0.94, 4.37, 2.02, 0.58, 15, WHITE, True, PP_ALIGN.CENTER)
    takeaway(slide, "SUBMIT: polarity RTL·TB·tested PPA  |  SEPARATE: CAV software evidence", GREEN)

    # Restamp headers last so the PowerPoint batch renderer cannot drop them
    # behind dense slide content.
    for finished_slide in list(prs.slides)[1:]:
        early = list(finished_slide.shapes)
        section_label = early[2].text
        title_label = early[3].text
        shape(finished_slide, 0.45, 0.14, 12.30, 0.68, BG, False)
        shape(finished_slide, 0.55, 0.26, 1.40, 0.34, CARD_2)
        text(finished_slide, section_label, 0.60, 0.30, 1.30, 0.22, 11, CYAN, True, PP_ALIGN.CENTER)
        text(finished_slide, title_label, 2.18, 0.18, 10.45, 0.58, 27, WHITE, True)
        shape(finished_slide, 0.55, 0.88, 12.18, 0.025, CYAN, False)

    prs.core_properties.title = "Cluster2 Steal-Buffer Polarity AER"
    prs.core_properties.subject = "AER bottleneck, architecture, quantitative improvement, RTL simulation, physical implementation, CAV extension"
    prs.core_properties.author = "AI-semi team"
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"WROTE {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
